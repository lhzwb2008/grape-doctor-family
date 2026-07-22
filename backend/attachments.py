"""临时附件解析：内存/临时目录处理，不写入业务 data 目录。"""

from __future__ import annotations

import io
import os
import re
import shutil
import subprocess
import tempfile
from pathlib import Path


MAX_TEXT_CHARS = 80000


def _clip(text: str) -> str:
    text = (text or "").strip()
    if len(text) > MAX_TEXT_CHARS:
        return text[:MAX_TEXT_CHARS] + "\n…(内容过长已截断)"
    return text


def _ext(name: str) -> str:
    return Path(name).suffix.lower()


def extract_plain_text(raw: bytes) -> str:
    for enc in ("utf-8", "utf-16", "gb18030", "latin-1"):
        try:
            return _clip(raw.decode(enc))
        except UnicodeDecodeError:
            continue
    return _clip(raw.decode("utf-8", errors="replace"))


def extract_pdf_text(raw: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(raw))
    parts: list[str] = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            t = page.extract_text() or ""
        except Exception:  # noqa: BLE001
            t = ""
        if t.strip():
            parts.append(f"--- 第 {i} 页 ---\n{t.strip()}")
    text = _clip("\n\n".join(parts))
    if not text:
        raise ValueError("未能从 PDF 提取到文字（可能是扫描件/图片版，请改传图片或粘贴文字）")
    return text


def extract_docx_text(raw: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(raw))
    parts = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    text = _clip("\n".join(parts))
    if not text:
        raise ValueError("未能从 Word 文档提取到文字")
    return text


def _extract_doc_via_cmd(raw: bytes, cmd: list[str]) -> str | None:
    """用系统工具解析旧版 .doc，文件只写在临时目录。"""
    binary = cmd[0]
    if not shutil.which(binary):
        return None
    with tempfile.TemporaryDirectory(prefix="grape-doc-") as tmp:
        path = Path(tmp) / "input.doc"
        path.write_bytes(raw)
        try:
            if binary in {"antiword", "catdoc"}:
                proc = subprocess.run(
                    [*cmd, str(path)],
                    capture_output=True,
                    timeout=45,
                    check=False,
                )
                out = (proc.stdout or b"").decode("utf-8", errors="replace").strip()
                if out:
                    return _clip(out)
            if binary in {"soffice", "libreoffice"}:
                proc = subprocess.run(
                    [
                        binary,
                        "--headless",
                        "--nologo",
                        "--nolockcheck",
                        "--nodefault",
                        "--nofirststartwizard",
                        "--convert-to",
                        "txt:Text",
                        "--outdir",
                        tmp,
                        str(path),
                    ],
                    capture_output=True,
                    timeout=90,
                    check=False,
                )
                txt_path = Path(tmp) / "input.txt"
                if txt_path.exists():
                    return _clip(txt_path.read_text(encoding="utf-8", errors="replace"))
                # 有些版本输出名不同
                for p in Path(tmp).glob("*.txt"):
                    return _clip(p.read_text(encoding="utf-8", errors="replace"))
                _ = proc
        except Exception:  # noqa: BLE001
            return None
    return None


def _extract_doc_ole_strings(raw: bytes) -> str:
    """无系统工具时的兜底：从 OLE 二进制里尽量捞可读文本。"""
    chunks: list[str] = []

    def _harvest(data: bytes) -> None:
        # UTF-16LE
        try:
            s = data.decode("utf-16le", errors="ignore")
        except Exception:  # noqa: BLE001
            s = ""
        for m in re.finditer(
            r"[\u4e00-\u9fffA-Za-z0-9][\u4e00-\u9fffA-Za-z0-9\s，。、；：！？“”‘’（）()\-_.，,]{3,}",
            s,
        ):
            chunks.append(m.group().strip())
        # ASCII
        for m in re.finditer(rb"[\x20-\x7e]{6,}", data):
            chunks.append(m.group().decode("ascii", errors="ignore").strip())

    try:
        import olefile

        if olefile.isOleFile(io.BytesIO(raw)):
            with olefile.OleFileIO(io.BytesIO(raw)) as ole:
                for stream_name in ("WordDocument", "1Table", "0Table"):
                    if ole.exists(stream_name):
                        _harvest(ole.openstream(stream_name).read())
        else:
            _harvest(raw)
    except Exception:  # noqa: BLE001
        _harvest(raw)

    seen: set[str] = set()
    uniq: list[str] = []
    for c in chunks:
        c = re.sub(r"\s+", " ", c).strip()
        if len(c) < 4 or c in seen:
            continue
        if c.count("?") > max(2, len(c) // 3):
            continue
        seen.add(c)
        uniq.append(c)
    text = _clip("\n".join(uniq))
    if not text:
        raise ValueError("未能从旧版 .doc 提取到可用文字")
    return text


def extract_doc_text(raw: bytes) -> str:
    for cmd in (
        ["antiword", "-m", "UTF-8.txt"],
        ["antiword"],
        ["catdoc", "-w"],
        ["soffice"],
        ["libreoffice"],
    ):
        text = _extract_doc_via_cmd(raw, cmd)
        if text:
            return text
    return _extract_doc_ole_strings(raw)


def extract_xlsx_text(raw: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"## 工作表：{ws.title}")
        row_n = 0
        for row in ws.iter_rows(values_only=True):
            vals = ["" if v is None else str(v).strip() for v in row]
            if not any(vals):
                continue
            parts.append(" | ".join(vals))
            row_n += 1
            if row_n >= 500:
                parts.append("…(该表行数过多已截断)")
                break
    wb.close()
    text = _clip("\n".join(parts))
    if not text:
        raise ValueError("Excel 内容为空")
    return text


def extract_xls_text(raw: bytes) -> str:
    import xlrd

    book = xlrd.open_workbook(file_contents=raw)
    parts: list[str] = []
    for sheet in book.sheets():
        parts.append(f"## 工作表：{sheet.name}")
        for r in range(min(sheet.nrows, 500)):
            vals = [str(sheet.cell_value(r, c)).strip() for c in range(sheet.ncols)]
            if any(vals):
                parts.append(" | ".join(vals))
        if sheet.nrows > 500:
            parts.append("…(该表行数过多已截断)")
    text = _clip("\n".join(parts))
    if not text:
        raise ValueError("Excel 内容为空")
    return text


def extract_pptx_text(raw: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(raw))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            parts.append(f"--- 第 {i} 页 ---\n" + "\n".join(texts))
    text = _clip("\n\n".join(parts))
    if not text:
        raise ValueError("未能从 PPT 提取到文字")
    return text


def extract_rtf_text(raw: bytes) -> str:
    from striprtf.striprtf import rtf_to_text

    text = _clip(rtf_to_text(raw.decode("latin-1", errors="ignore")))
    if not text:
        raise ValueError("未能从 RTF 提取到文字")
    return text


def extract_attachment_text(name: str, mime: str, raw: bytes) -> tuple[str, str]:
    """
    返回 (类型标签, 正文)。
    失败抛 ValueError。
    """
    mime = (mime or "").split(";")[0].strip().lower()
    name = name or "attachment"
    ext = _ext(name)

    # 文本
    if (
        mime.startswith("text/")
        or mime in {"application/json", "application/xml", "text/xml", "text/csv", "text/markdown"}
        or ext in {".txt", ".md", ".csv", ".json", ".log", ".py", ".js", ".ts", ".html", ".css", ".xml", ".yaml", ".yml"}
    ):
        return "文本", extract_plain_text(raw)

    # PDF
    if mime in {"application/pdf", "application/x-pdf"} or ext == ".pdf":
        return "PDF", extract_pdf_text(raw)

    # Word
    if (
        mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        or ext == ".docx"
    ):
        return "Word", extract_docx_text(raw)

    if mime == "application/msword" or ext == ".doc":
        return "Word(.doc)", extract_doc_text(raw)

    # Excel
    if (
        mime == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        or ext == ".xlsx"
    ):
        return "Excel", extract_xlsx_text(raw)

    if mime in {"application/vnd.ms-excel", "application/excel"} or ext == ".xls":
        return "Excel(.xls)", extract_xls_text(raw)

    # PPT
    if (
        mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        or ext == ".pptx"
    ):
        return "PPT", extract_pptx_text(raw)

    # RTF
    if mime == "application/rtf" or ext == ".rtf":
        return "RTF", extract_rtf_text(raw)

    raise ValueError(f"暂不支持该格式（{ext or mime or 'unknown'}）")
